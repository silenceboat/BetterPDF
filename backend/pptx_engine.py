"""
PPTX Engine for rendering PowerPoint presentations as paginated PNG images.
"""

import base64
import platform
from io import BytesIO
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from PIL import Image, ImageDraw, ImageFont


# Standard slide canvas dimensions in points (10" x 7.5" at 72 dpi)
SLIDE_WIDTH = 720
SLIDE_HEIGHT = 540

# Default font size when none is specified
DEFAULT_FONT_SIZE = 18


class PptxEngine:
    """
    PPTX rendering engine.

    Parses a .pptx file with python-pptx and renders each slide
    as a base64-encoded PNG image using Pillow.
    """

    SLIDE_WIDTH = SLIDE_WIDTH
    SLIDE_HEIGHT = SLIDE_HEIGHT

    def __init__(self, file_path: str):
        self.file_path = file_path
        self._font_cache: Dict[Tuple[int, bool], ImageFont.FreeTypeFont] = {}
        self._render_cache: Dict[Tuple[int, float], str] = {}

        try:
            import pptx
        except ImportError:
            raise ImportError(
                "python-pptx is required for PPT support. "
                "Install it with: pip install python-pptx"
            )

        self._prs = pptx.Presentation(file_path)
        self.page_count = len(self._prs.slides)

        # Slide dimensions in EMU (English Metric Units)
        self._slide_width_emu: int = int(self._prs.slide_width or 9144000)
        self._slide_height_emu: int = int(self._prs.slide_height or 6858000)

        self._font_normal_paths, self._font_bold_paths = self._get_font_paths()

    # ------------------------------------------------------------------
    # Font helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _get_font_paths() -> Tuple[List[str], List[str]]:
        """Return (normal_paths, bold_paths) for the current platform."""
        system = platform.system()
        if system == "Windows":
            return (
                ["C:/Windows/Fonts/msyh.ttc", "C:/Windows/Fonts/arial.ttf"],
                ["C:/Windows/Fonts/msyhbd.ttc", "C:/Windows/Fonts/arialbd.ttf"],
            )
        elif system == "Linux":
            return (
                [
                    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                    "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
                ],
                [
                    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
                    "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
                ],
            )
        else:  # macOS / other
            return (
                ["/System/Library/Fonts/PingFang.ttc"],
                ["/System/Library/Fonts/PingFang.ttc"],
            )

    def _load_font(self, size: int, bold: bool) -> ImageFont.FreeTypeFont:
        """Load a font at the given size, using cache."""
        key = (size, bold)
        if key in self._font_cache:
            return self._font_cache[key]

        paths = self._font_bold_paths if bold else self._font_normal_paths
        font = None
        for p in paths:
            try:
                font = ImageFont.truetype(p, size)
                break
            except (OSError, IOError):
                continue

        if font is None and bold:
            # Fallback to normal font for bold
            for p in self._font_normal_paths:
                try:
                    font = ImageFont.truetype(p, size)
                    break
                except (OSError, IOError):
                    continue

        if font is None:
            font = ImageFont.load_default()

        self._font_cache[key] = font
        return font

    # ------------------------------------------------------------------
    # Coordinate helpers
    # ------------------------------------------------------------------

    def _shape_rect(self, shape, scale: float) -> Tuple[int, int, int, int]:
        """Return (x, y, w, h) in pixels for a shape at the given scale."""
        sx = SLIDE_WIDTH / self._slide_width_emu * scale
        sy = SLIDE_HEIGHT / self._slide_height_emu * scale
        x = int((shape.left or 0) * sx)
        y = int((shape.top or 0) * sy)
        w = int((shape.width or 0) * sx)
        h = int((shape.height or 0) * sy)
        return x, y, w, h

    # ------------------------------------------------------------------
    # Background extraction
    # ------------------------------------------------------------------

    def _get_slide_bg_color(self, slide) -> Tuple[int, int, int]:
        """Extract background fill color from slide/layout/master. Defaults to white."""
        for source in (slide, slide.slide_layout, slide.slide_layout.slide_master):
            try:
                fill = source.background.fill
                if fill.type is not None:
                    rgb = fill.fore_color.rgb
                    return (rgb.red, rgb.green, rgb.blue)
            except Exception:
                continue
        return (255, 255, 255)

    # ------------------------------------------------------------------
    # Text helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _resolve_font_size(run_font, para_font) -> int:
        """Pick font size in points from run or paragraph font (fallback to default)."""
        for font in (run_font, para_font):
            try:
                if font and font.size:
                    pt = font.size.pt
                    if pt and pt > 0:
                        return max(6, int(pt))
            except Exception:
                continue
        return DEFAULT_FONT_SIZE

    @staticmethod
    def _resolve_bold(run_font, para_font) -> bool:
        """Determine if text should be bold."""
        for font in (run_font, para_font):
            try:
                if font and font.bold is not None:
                    return bool(font.bold)
            except Exception:
                continue
        return False

    @staticmethod
    def _resolve_color(run_font, para_font) -> Tuple[int, int, int]:
        """Extract text color (R, G, B). Defaults to black."""
        for font in (run_font, para_font):
            try:
                if font and font.color and font.color.type is not None:
                    rgb = font.color.rgb
                    return (rgb.red, rgb.green, rgb.blue)
            except Exception:
                continue
        return (0, 0, 0)

    @staticmethod
    def _wrap_text(text: str, font: ImageFont.FreeTypeFont, max_width: float) -> List[str]:
        """Greedy character-level wrapping that handles CJK."""
        if not text:
            return [""]
        lines: List[str] = []
        current_line = ""
        current_width = 0.0
        for char in text:
            try:
                cw = font.getlength(char)
            except AttributeError:
                bbox = font.getbbox(char)
                cw = (bbox[2] - bbox[0]) if bbox else 6.0
            if current_width + cw <= max_width:
                current_line += char
                current_width += cw
            else:
                if current_line:
                    lines.append(current_line)
                current_line = char
                current_width = cw
        if current_line:
            lines.append(current_line)
        return lines or [""]

    # ------------------------------------------------------------------
    # Shape rendering
    # ------------------------------------------------------------------

    def _render_picture(self, img: Image.Image, shape, x: int, y: int, w: int, h: int):
        """Embed a picture shape into the canvas."""
        try:
            blob = shape.image.blob
            pic = Image.open(BytesIO(blob)).convert("RGBA")
            if w > 0 and h > 0:
                pic = pic.resize((w, h), Image.LANCZOS)
            # Paste with alpha mask if available
            if pic.mode == "RGBA":
                img.paste(pic, (x, y), pic)
            else:
                img.paste(pic, (x, y))
        except Exception:
            pass

    def _render_text_frame(
        self,
        draw: ImageDraw.Draw,
        text_frame,
        x: int,
        y: int,
        w: int,
        h: int,
        scale: float,
    ):
        """Render all paragraphs of a text frame onto the canvas."""
        padding = max(2, int(4 * scale))
        max_text_w = max(20, w - 2 * padding) if w > 0 else int(SLIDE_WIDTH * scale * 0.9)
        cursor_y = y + padding

        para_font = text_frame.paragraphs[0].font if text_frame.paragraphs else None

        for para in text_frame.paragraphs:
            para_font_obj = para.font

            # Collect runs; treat empty paragraph as blank line
            if not para.runs and not para.text:
                line_h = int(DEFAULT_FONT_SIZE * scale * 1.3)
                cursor_y += line_h
                continue

            # Use first run for style (representative)
            first_run_font = para.runs[0].font if para.runs else None
            font_size_pt = self._resolve_font_size(first_run_font, para_font_obj)
            bold = self._resolve_bold(first_run_font, para_font_obj)
            color = self._resolve_color(first_run_font, para_font_obj)

            scaled_size = max(6, int(font_size_pt * scale))
            font = self._load_font(scaled_size, bold)
            line_h = int(scaled_size * 1.35)

            text = para.text
            if not text:
                cursor_y += line_h
                continue

            wrapped = self._wrap_text(text, font, max_text_w)
            for line in wrapped:
                bottom_limit = y + h if h > 0 else int(SLIDE_HEIGHT * scale)
                if cursor_y + line_h > bottom_limit:
                    return  # clip to shape
                draw.text((x + padding, cursor_y), line, fill=color, font=font)
                cursor_y += line_h

    def _render_shape(
        self,
        draw: ImageDraw.Draw,
        img: Image.Image,
        shape,
        scale: float,
    ):
        """Render a single slide shape."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore

        x, y, w, h = self._shape_rect(shape, scale)

        # Pictures
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            self._render_picture(img, shape, x, y, w, h)
            return

        # Optionally draw shape background fill
        try:
            fill = shape.fill
            if fill.type is not None:
                rgb = fill.fore_color.rgb
                fill_color = (rgb.red, rgb.green, rgb.blue)
                if w > 0 and h > 0:
                    draw.rectangle([x, y, x + w, y + h], fill=fill_color)
        except Exception:
            pass

        # Text frames
        if shape.has_text_frame:
            self._render_text_frame(draw, shape.text_frame, x, y, w, h, scale)

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def get_metadata(self) -> dict:
        """Get presentation metadata."""
        title = ""
        author = ""
        try:
            props = self._prs.core_properties
            title = props.title or ""
            author = props.author or ""
        except Exception:
            pass
        return {
            "file_name": Path(self.file_path).name,
            "page_count": self.page_count,
            "title": title or Path(self.file_path).stem,
            "author": author,
            "subject": "PowerPoint Presentation",
        }

    def render_page(self, page_num: int, zoom: float = 1.0) -> str:
        """
        Render a slide to a base64-encoded PNG image.

        Args:
            page_num: 1-based slide number
            zoom: Zoom factor

        Returns:
            Base64-encoded PNG string (without data URI prefix)
        """
        cache_key = (page_num, round(zoom, 2))
        if cache_key in self._render_cache:
            return self._render_cache[cache_key]

        if page_num < 1 or page_num > self.page_count:
            raise ValueError(f"Invalid page number: {page_num}")

        slide = self._prs.slides[page_num - 1]
        img_w = max(1, int(SLIDE_WIDTH * zoom))
        img_h = max(1, int(SLIDE_HEIGHT * zoom))

        bg = self._get_slide_bg_color(slide)
        img = Image.new("RGB", (img_w, img_h), color=bg)
        draw = ImageDraw.Draw(img)

        for shape in slide.shapes:
            try:
                self._render_shape(draw, img, shape, zoom)
            except Exception:
                pass  # Skip shapes that can't be rendered

        buffer = BytesIO()
        img.save(buffer, format="PNG", optimize=True)
        result = base64.b64encode(buffer.getvalue()).decode()

        self._render_cache[cache_key] = result
        return result

    def extract_text(self, page_num: int, rect: Optional[dict] = None) -> str:
        """Extract all text from a slide."""
        if page_num < 1 or page_num > self.page_count:
            raise ValueError(f"Invalid page number: {page_num}")

        slide = self._prs.slides[page_num - 1]
        parts: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
        return "\n".join(parts)

    def get_page_size(self, page_num: int) -> Tuple[float, float]:
        """Return slide canvas size in points."""
        if page_num < 1 or page_num > self.page_count:
            raise ValueError(f"Invalid page number: {page_num}")
        return (float(SLIDE_WIDTH), float(SLIDE_HEIGHT))

    def search_text(self, query: str, page_num: Optional[int] = None) -> list:
        """
        Search for text across slides.

        Returns list of {page, rect: {x1, y1, x2, y2}} dicts.
        """
        results = []
        pages = [page_num] if page_num else range(1, self.page_count + 1)
        query_lower = query.lower()

        for pn in pages:
            slide = self._prs.slides[pn - 1]
            sx = SLIDE_WIDTH / self._slide_width_emu
            sy = SLIDE_HEIGHT / self._slide_height_emu

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    if query_lower not in para.text.lower():
                        continue
                    x1 = float((shape.left or 0) * sx)
                    y1 = float((shape.top or 0) * sy)
                    x2 = float(x1 + (shape.width or 0) * sx)
                    y2 = float(y1 + (shape.height or 0) * sy)
                    results.append({
                        "page": pn,
                        "rect": {"x1": x1, "y1": y1, "x2": x2, "y2": y2},
                    })

        return results

    def close(self):
        """Free resources."""
        self._render_cache.clear()
        self._font_cache.clear()

    def __enter__(self):
        return self

    def __exit__(self, exc_type, exc_val, exc_tb):
        self.close()
